#!/usr/bin/env python3
"""ROTA 组会 PPT — 技术线为主，说人话，不搞背景铺垫。"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "ROTA_组会汇报_2026-08-19.pptx"
DL = ROOT / "apps" / "rota" / "static" / "downloads" / "ROTA_group_meeting.pptx"

GREEN = RGBColor(0x1A, 0x3A, 0x2E)
BODY = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x55, 0x55, 0x55)


def _title(slide, text: str, size=30):
    slide.shapes.title.text = text
    for p in slide.shapes.title.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = True
            r.font.color.rgb = GREEN


def _bullets(prs, title: str, items: list[str], sub: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _title(slide, title)
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, t in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = t
        p.font.size = Pt(22)
        p.font.color.rgb = BODY
        p.space_after = Pt(12)
    if sub:
        p = tf.add_paragraph()
        p.text = sub
        p.font.size = Pt(17)
        p.font.color.rgb = MUTED
        p.space_before = Pt(16)


def _flow(prs, title: str, lines: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    tx.text_frame.text = title
    for p in tx.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(28)
            r.font.bold = True
            r.font.color.rgb = GREEN
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.8))
    tf = box.text_frame
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "PingFang SC"
        p.font.size = Pt(20 if line.startswith(" ") or "→" in line else 22)
        p.font.color.rgb = BODY if not line.startswith("【") else MUTED
        p.space_after = Pt(6)


def _two(prs, title: str, a: list[str], b: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    tx.text_frame.text = title
    for p in tx.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(28)
            r.font.bold = True
            r.font.color.rgb = GREEN

    def col(x, items):
        b = slide.shapes.add_textbox(Inches(x), Inches(1.2), Inches(4.3), Inches(5.8))
        tf = b.text_frame
        tf.word_wrap = True
        for i, t in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = t
            p.font.size = Pt(20)
            p.space_after = Pt(10)

    col(0.5, a)
    col(5.2, b)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 封面
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "ROTA 技术线进度"
    s.placeholders[1].text = "单目演示 · 2026-08-19"
    _title(s, "ROTA 技术线进度", 34)

    _bullets(
        prs,
        "两条线在并行",
        [
            "单目线（现在在推）：一条视频 → 2D → 约束 3D → 网页能看。T014 侧拍已跑通。",
            "多机线（金标准）：三机标定 + 同步 → 三角化。T012 音频同步失败，等视觉标定出 cameras.yaml。",
            "老师要的 monocular 先走单目；多机出来再用来验单目靠不靠谱。",
        ],
    )

    _flow(
        prs,
        "单目 pipeline（就这一条链）",
        [
            "视频 MP4",
            "  → 抽帧（10 fps）",
            "  → Sapiens2：2D 关键点，锁运动员",
            "  → MotionBERT：给一个深度初值（弱，不直接展示）",
            "  → 标 5 个车点：定车体坐标（后花鼓→前花鼓 = +X）",
            "  → bike_geometry v5：骑姿硬约束 + 踏频驱动",
            "  → joints_constrained.json → ROTA 网页 3D",
        ],
    )

    _two(
        prs,
        "3D 里什么是视频给的 / 什么是模型补的",
        [
            "来自视频",
            "· 2D 关节位置（Sapiens）",
            "· 踏频 / 曲柄相位（脚踝轨迹）",
            "· 可见侧肩髋的一点抖动",
            "",
            "不来自视频",
            "· 左膝左肘（右侧机位看不见）",
            "· 车架真实尺寸（用标准公路车模板）",
            "· 坐鞍、扶把、反相脚踏（硬约束）",
        ],
        [
            "半合成",
            "· 躯干左右晃：跟踏频同步，",
            "  幅度一部分按骑行常识加",
            "· MotionBERT 只贡献深度比例",
            "",
            "一句话",
            "不是多目重建；是",
            "「2D + 人会怎么骑车」",
            "的可视化。",
        ],
    )

    _bullets(
        prs,
        "bike_geometry 在干什么（核心模块）",
        [
            "手锁把套、脚锁踏板、左右脚踏差 180°，人坐鞍上——自行车当闭链机构解。",
            "膝、肘：两连杆 IK；右侧机位时左肢锁 +Z，不跟被压扁的 2D 走。",
            "v5 改动：肩髋按 sin(曲柄角) 左右晃 + 少量 2D 残差；蒙皮换 cyclist.glb。",
            "求解器：cycling_hard_kinematics_v5_upper_sway",
        ],
        sub="T014：80 帧 · ~105 rpm · 重投影 ~189 px（说明没贴回原图，别当测量）",
    )

    _bullets(
        prs,
        "现在能演示什么",
        [
            "ROTA 网页：左原视频、右 3D，时间轴同步，可标 5 车点重算。",
            "公网：…seetacloud.com:8443/app（实例要开机）",
            "离线：T014 有 side_by_side_2d_3d.mp4、joints_constrained.json",
        ],
    )

    _bullets(
        prs,
        "还没接上的",
        [
            "上传新视频 → 自动跑 Sapiens/MotionBERT：API 还是占位，换片要离线重跑脚本。",
            "多机：T012 棋盘格标定 → cameras.yaml → T011 最小三角化验证。",
            "apps/rota/ 大块代码未 commit。",
        ],
    )

    _bullets(
        prs,
        "下一步（按优先级）",
        [
            "1. /api/jobs 接真实批处理（上传就能跑）",
            "2. T012 视觉标定，出 cameras.yaml",
            "3. 同一 Trial 上单目 vs 多机比膝角/踏频",
        ],
    )

    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "谢谢"
    s.placeholders[1].text = "可现场开 /app 演示"

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    DL.parent.mkdir(parents=True, exist_ok=True)
    import shutil
    shutil.copy2(OUT, DL)
    return OUT


if __name__ == "__main__":
    p = build()
    print(f"wrote: {p}")
